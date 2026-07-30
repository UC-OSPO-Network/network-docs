"""
Apply UC OSPO branding to the Quarto reference PPTX.

Colors from ucospo-theme.scss:
  navy   #1f335e   headings, title slide bg
  blue   #003cb3   links, accent
  orange #ff9100   underlines, accents
  teal   #00a291   accent
  text   #222832   body
  bg     #ffffff   slide background
  muted  #f0f4fa   light fills
"""
import os
import zipfile
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt, Emu
from pptx.oxml.ns import qn

INPUT = os.path.join(os.path.dirname(__file__), "ucospo-reference.pptx")

# ── Step 1: Patch theme XML (colors + fonts) via zip ─────────────────────────

files = {}
with zipfile.ZipFile(INPUT, "r") as z:
    for name in z.namelist():
        files[name] = z.read(name)

theme_xml = files["ppt/theme/theme1.xml"].decode("utf-8")

replacements = [
    # Color scheme
    ('"1F497D"', '"1f335e"'),   # dk2  → navy
    ('"EEECE1"', '"f0f4fa"'),   # lt2  → muted
    ('"4F81BD"', '"003cb3"'),   # accent1 → blue
    ('"C0504D"', '"ff9100"'),   # accent2 → orange
    ('"9BBB59"', '"00a291"'),   # accent3 → teal
    ('"8064A2"', '"1f335e"'),   # accent4 → navy
    ('"4BACC6"', '"f0f4fa"'),   # accent5 → muted
    ('"F79646"', '"222832"'),   # accent6 → text
    ('"0000FF"', '"003cb3"'),   # hlink   → blue
    ('"800080"', '"ff9100"'),   # folHlink → orange
    # Fonts
    ('typeface="Calibri"', 'typeface="Arial"'),
]

for old, new in replacements:
    theme_xml = theme_xml.replace(old, new)

files["ppt/theme/theme1.xml"] = theme_xml.encode("utf-8")

tmp = INPUT + ".tmp"
with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as z:
    for name, data in files.items():
        z.writestr(name, data)
os.replace(tmp, INPUT)
print("Step 1 done: theme colors and fonts patched")


# ── Step 2: Slide layout styling via python-pptx ─────────────────────────────

prs = Presentation(INPUT)
master = prs.slide_masters[0]

NAVY   = RGBColor(0x1f, 0x33, 0x5e)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xFF, 0x91, 0x00)
TEXT   = RGBColor(0x22, 0x28, 0x32)
MUTED  = RGBColor(0xF0, 0xF4, 0xFA)


def set_solid_background(slide_or_layout, color: RGBColor):
    """Set a solid fill background on a slide/layout via XML."""
    sp_tree = slide_or_layout.shapes._spTree
    # Remove existing bg element if present
    for bg in sp_tree.getparent().findall(qn("p:bg")):
        sp_tree.getparent().remove(bg)

    # Build: <p:bg><p:bgPr><a:solidFill><a:srgbClr val="..."/></a:solidFill>...
    bg_xml = (
        f'<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<p:bgPr>'
        f'<a:solidFill><a:srgbClr val="{str(color).upper()}"/></a:solidFill>'
        f'<a:effectLst/>'
        f'</p:bgPr>'
        f'</p:bg>'
    )
    bg_el = etree.fromstring(bg_xml)
    # Insert before spTree
    cSld = sp_tree.getparent()
    cSld.insert(list(cSld).index(sp_tree), bg_el)


def set_placeholder_font(placeholder, color: RGBColor, bold=False, size_pt=None):
    """Set text color in a layout placeholder via lstStyle levels and existing runs.

    Pandoc inherits formatting from lstStyle, not from the sample run text in
    the reference PPTX, so we must set color at every level in lstStyle.
    """
    hex_color = f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"
    fill_xml = (
        f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:srgbClr val="{hex_color}"/>'
        f'</a:solidFill>'
    )

    txBody = placeholder._element.find(qn("p:txBody"))
    if txBody is None:
        return

    # Set/replace lstStyle with explicit color at all 9 levels
    lst = txBody.find(qn("a:lstStyle"))
    if lst is None:
        lst = etree.SubElement(txBody, qn("a:lstStyle"))
    else:
        lst.clear()

    level_tags = [
        "a:lvl1pPr", "a:lvl2pPr", "a:lvl3pPr", "a:lvl4pPr",
        "a:lvl5pPr", "a:lvl6pPr", "a:lvl7pPr", "a:lvl8pPr", "a:lvl9pPr",
    ]
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    for tag in level_tags:
        lvl = etree.SubElement(lst, f"{{{ns}}}{tag.split(':')[1]}")
        defRPr = etree.SubElement(lvl, f"{{{ns}}}defRPr")
        defRPr.append(etree.fromstring(fill_xml))
        if bold:
            defRPr.set("b", "1")
        if size_pt:
            defRPr.set("sz", str(int(size_pt * 100)))

    # Also fix any explicit runs already in the placeholder (sample text)
    for rPr in placeholder._element.iter(qn("a:rPr")):
        for sf in rPr.findall(qn("a:solidFill")):
            rPr.remove(sf)
        for sc in rPr.findall(qn("a:schemeClr")):
            rPr.remove(sc)
        rPr.insert(0, etree.fromstring(fill_xml))


def add_orange_rule(layout, placeholder_idx=0):
    """Add a 3pt orange horizontal rule below the title placeholder via XML."""
    ph = None
    for shape in layout.placeholders:
        if shape.placeholder_format.idx == placeholder_idx:
            ph = shape
            break
    if ph is None:
        return

    left   = ph.left
    top    = ph.top + ph.height + 76200        # ~0.08" below title
    width  = ph.width
    height = 38100                              # ~3pt / 0.04"

    sp_xml = (
        f'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        f' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        f'<p:nvSpPr>'
        f'<p:cNvPr id="100" name="OrangeRule"/>'
        f'<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
        f'<p:nvPr/>'
        f'</p:nvSpPr>'
        f'<p:spPr>'
        f'<a:xfrm><a:off x="{left}" y="{top}"/><a:ext cx="{width}" cy="{height}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f'<a:solidFill><a:srgbClr val="FF9100"/></a:solidFill>'
        f'<a:ln><a:noFill/></a:ln>'
        f'</p:spPr>'
        f'<p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>'
        f'</p:sp>'
    )
    sp_el = etree.fromstring(sp_xml)
    layout.shapes._spTree.append(sp_el)


# Layout 0: Title Slide — navy background, white text
title_layout = master.slide_layouts[0]
set_solid_background(title_layout, NAVY)
for ph in title_layout.placeholders:
    set_placeholder_font(ph, WHITE)
print("Step 2a done: title slide layout styled")

# Layout 1: Title and Content — white bg, navy heading, orange rule
BODY_FONT_PT = 16
content_layout = master.slide_layouts[1]
set_solid_background(content_layout, WHITE)
for ph in content_layout.placeholders:
    if ph.placeholder_format.idx == 0:   # title
        set_placeholder_font(ph, NAVY, bold=True)
    else:
        set_placeholder_font(ph, TEXT, size_pt=BODY_FONT_PT)
add_orange_rule(content_layout, placeholder_idx=0)
print("Step 2b done: content layout styled")

# Layout 2: Section Header — navy background, white text
section_layout = master.slide_layouts[2]
set_solid_background(section_layout, NAVY)
for ph in section_layout.placeholders:
    set_placeholder_font(ph, WHITE)
print("Step 2c done: section header layout styled")

# Layouts 3+ : white background, navy title, text body, orange rule on title
# Layout 3 = Two Content, 4 = Comparison, 5 = Title Only, etc.
BODY_FONT_PT = 16  # reduce from default ~24pt to prevent overflow

for i, layout in enumerate(master.slide_layouts):
    if i in (0, 1, 2):
        continue
    set_solid_background(layout, WHITE)
    for ph in layout.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            set_placeholder_font(ph, NAVY, bold=True)
        else:
            set_placeholder_font(ph, TEXT, size_pt=BODY_FONT_PT)
    # Add orange rule to any layout that has a title placeholder (idx 0)
    has_title = any(p.placeholder_format.idx == 0 for p in layout.placeholders)
    if has_title:
        add_orange_rule(layout, placeholder_idx=0)
print("Step 2d done: remaining layouts styled with orange rules and reduced body font")

prs.save(INPUT)
print(f"\nDone. Saved to: {INPUT}")
